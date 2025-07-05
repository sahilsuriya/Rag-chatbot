import pdfplumber, pandas as pd, os
from pptx import Presentation
from docx import Document

def parse_txt(file_path):
    with open(file_path, "r", encoding="utf-8") as f:
        return f.read()

def parse_pdf(file_path):
    with pdfplumber.open(file_path) as pdf:
        return "\n".join([page.extract_text() or '' for page in pdf.pages])

def parse_docx(file_path):
    doc = Document(file_path)
    return "\n".join([p.text for p in doc.paragraphs])

def parse_pptx(file_path):
    prs = Presentation(file_path)
    return "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])

def parse_csv(file_path):
    df = pd.read_csv(file_path)
    return df.to_string()

def parse_document(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".pdf": return parse_pdf(file_path)
    if ext == ".txt" or ext == ".md": return parse_txt(file_path)
    if ext == ".docx": return parse_docx(file_path)
    if ext == ".pptx": return parse_pptx(file_path)
    if ext == ".csv": return parse_csv(file_path)
    return ""
